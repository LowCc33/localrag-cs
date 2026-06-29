#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文件解析器模块
支持多种格式文件的内容提取：.md, .txt, .pdf, .docx

职责：
1. 根据文件后缀自动选择解析器
2. 处理不同编码格式
3. 提取纯文本内容，去除格式标记
4. 异常处理：解析失败时记录原因，不中断整体流程
5. 支持中文字符编码

设计原则：
- 最小依赖：只引入必要的解析库
- 错误容忍：单个文件解析失败不影响其他文件
- 性能优化：大文件分块读取，避免内存溢出
- 编码兼容：自动检测和处理UTF-8、GBK等编码
"""

import os
import re
import logging
from pathlib import Path
from typing import Optional, List
from abc import ABC, abstractmethod

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


class BaseParser(ABC):
    """文件解析器基类"""
    
    def __init__(self):
        self.supported_extensions: List[str] = []
    
    @abstractmethod
    def parse(self, file_path: Path) -> Optional[str]:
        """
        解析文件内容
        
        Args:
            file_path: 文件路径对象
            
        Returns:
            解析出的纯文本内容，解析失败返回None
        """
        pass
    
    def extract_title(self, file_path: Path, content: Optional[str] = None) -> str:
        """
        提取文档标题
        
        默认使用文件名（不含扩展名）作为标题。
        子类可覆盖此方法以提取更准确的标题（如 Markdown 的第一个 # 标题）。
        
        Args:
            file_path: 文件路径对象
            content: 已解析的文本内容（可选，部分解析器需要从内容中提取标题）
            
        Returns:
            文档标题字符串
        """
        return file_path.stem
        """
        解析文件内容
        
        Args:
            file_path: 文件路径对象
            
        Returns:
            解析出的纯文本内容，解析失败返回None
            
        Raises:
            不抛出异常，解析失败返回None并记录日志
        """
        pass
    
    def can_parse(self, file_path: Path) -> bool:
        """检查是否支持解析此文件"""
        return file_path.suffix.lower() in self.supported_extensions


class TextParser(BaseParser):
    """纯文本文件解析器（.txt）"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.txt']
    
    def parse(self, file_path: Path) -> Optional[str]:
        """
        解析纯文本文件
        
        支持多种编码格式：
        - UTF-8 (带/不带BOM)
        - GBK/GB2312 (中文编码)
        - 自动检测编码，尝试多种编码格式
        
        Args:
            file_path: 文本文件路径
            
        Returns:
            文件文本内容，解析失败返回None
        """
        try:
            # 尝试多种编码格式
            encodings = ['utf-8', 'utf-8-sig', 'gbk', 'gb2312', 'latin-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    
                    # 验证是否成功读取（非空且包含可打印字符）
                    if content and any(c.isprintable() for c in content):
                        logger.debug(f"成功读取文件 {file_path.name}，编码: {encoding}")
                        return content
                        
                except UnicodeDecodeError:
                    continue
            
            # 所有编码都失败
            logger.error(f"无法解码文件 {file_path.name}，尝试的编码: {encodings}")
            return None
            
        except Exception as e:
            logger.error(f"解析文本文件失败 {file_path}: {e}")
            return None


class MarkdownParser(BaseParser):
    """Markdown文件解析器（.md）"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.md', '.markdown']
    
    def extract_title(self, file_path: Path, content: Optional[str] = None) -> str:
        """
        从 Markdown 内容中提取第一个 # 标题
        
        如果找不到 # 标题，回退到文件名。
        
        Args:
            file_path: 文件路径对象
            content: 已解析的 Markdown 文本内容
            
        Returns:
            文档标题
        """
        if content:
            # 匹配第一个 # 标题（支持 ## 和 ### 等，但取第一个 # 开头的行）
            match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
            if match:
                return match.group(1).strip()
        # 回退到文件名
        return file_path.stem
    
    def parse(self, file_path: Path) -> Optional[str]:
        """
        解析Markdown文件
        
        处理策略：
        1. 去除Markdown格式标记（#标题、*列表、**加粗**等）
        2. 保留代码块内容但去除```标记
        3. 保留链接文本但去除[链接](url)格式
        4. 保留图片alt文本但去除![alt](url)格式
        
        Args:
            file_path: Markdown文件路径
            
        Returns:
            清理后的纯文本内容
        """
        try:
            # 先按文本文件读取
            text_parser = TextParser()
            content = text_parser.parse(file_path)
            
            if not content:
                return None
            
            # 清理Markdown格式标记
            cleaned = self._clean_markdown(content)
            return cleaned
            
        except Exception as e:
            logger.error(f"解析Markdown文件失败 {file_path}: {e}")
            return None
    
    def _clean_markdown(self, content: str) -> str:
        """
        清理Markdown格式标记
        
        Args:
            content: 原始Markdown内容
            
        Returns:
            清理后的纯文本
        """
        # 移除标题标记（# ## ###等）
        content = re.sub(r'^#+\s+', '', content, flags=re.MULTILINE)
        
        # 移除粗体、斜体标记（**text**、*text*）
        content = re.sub(r'\*\*(.*?)\*\*', r'\1', content)
        content = re.sub(r'\*(.*?)\*', r'\1', content)
        
        # 移除删除线标记（~~text~~）
        content = re.sub(r'~~(.*?)~~', r'\1', content)
        
        # 处理行内代码标记（`code`）
        content = re.sub(r'`(.*?)`', r'\1', content)
        
        # 处理链接：[text](url) -> text
        content = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', content)
        
        # 处理图片：![alt](url) -> alt
        content = re.sub(r'!\[([^\]]*)\]\([^)]+\)', r'\1', content)
        
        # 移除代码块标记（```language ... ```）
        content = re.sub(r'```[\s\S]*?```', '', content)
        
        # 移除无序列表标记（-、*、+）
        content = re.sub(r'^[-*+]\s+', '', content, flags=re.MULTILINE)
        
        # 移除有序列表标记（1.、2.等）
        content = re.sub(r'^\d+\.\s+', '', content, flags=re.MULTILINE)
        
        # 移除引用标记（>）
        content = re.sub(r'^>\s+', '', content, flags=re.MULTILINE)
        
        # 移除水平线（---、***）
        content = re.sub(r'^[-*]{3,}\s*$', '', content, flags=re.MULTILINE)
        
        # 合并多个空白行
        content = re.sub(r'\n\s*\n\s*\n', '\n\n', content)
        
        return content.strip()


class PDFParser(BaseParser):
    """PDF文件解析器"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.pdf']
        self._has_pypdf2 = self._check_pypdf2_availability()
    
    def _check_pypdf2_availability(self) -> bool:
        """检查PyPDF2是否可用"""
        try:
            import importlib.util
            spec = importlib.util.find_spec("PyPDF2")
            return spec is not None
        except Exception:
            logger.warning("PyPDF2库未安装，PDF解析功能受限")
            return False
    
    def parse(self, file_path: Path) -> Optional[str]:
        """
        解析PDF文件
        
        使用PyPDF2库提取文本内容（如果可用）
        否则尝试使用命令行工具pdftotext
        
        Args:
            file_path: PDF文件路径
            
        Returns:
            PDF文本内容，解析失败返回None
        """
        if self._has_pypdf2:
            return self._parse_with_pypdf2(file_path)
        else:
            return self._parse_with_pdftotext(file_path)
    
    def _parse_with_pypdf2(self, file_path: Path) -> Optional[str]:
        """使用PyPDF2解析PDF"""
        try:
            import PyPDF2
            
            content_parts = []
            
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                
                # 检查PDF是否加密
                if pdf_reader.is_encrypted:
                    logger.warning(f"PDF文件已加密，跳过: {file_path.name}")
                    return None
                
                # 逐页提取文本
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    try:
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            content_parts.append(page_text.strip())
                            
                            # 每10页记录一次进度
                            if page_num % 10 == 0:
                                logger.debug(f"PDF解析进度: {page_num}/{len(pdf_reader.pages)} 页")
                    except Exception as e:
                        logger.warning(f"PDF第{page_num}页解析失败: {e}")
                        continue
            
            if not content_parts:
                logger.warning(f"PDF文件无文本内容: {file_path.name}")
                return None
            
            # 合并所有页面内容
            full_content = '\n\n'.join(content_parts)
            logger.info(f"PDF解析完成: {file_path.name}，共{len(pdf_reader.pages)}页")
            
            return full_content
            
        except Exception as e:
            logger.error(f"PyPDF2解析PDF文件失败 {file_path}: {e}")
            # 尝试使用pdftotext作为备选
            return self._parse_with_pdftotext(file_path)
    
    def _parse_with_pdftotext(self, file_path: Path) -> Optional[str]:
        """使用命令行工具pdftotext解析PDF"""
        try:
            import subprocess
            import tempfile
            
            # 检查pdftotext是否可用
            result = subprocess.run(['which', 'pdftotext'], 
                                  capture_output=True, text=True)
            if result.returncode != 0:
                logger.error("pdftotext命令不可用，请安装poppler-utils")
                return None
            
            # 创建临时文件用于输出
            with tempfile.NamedTemporaryFile(mode='w+', suffix='.txt', delete=False) as tmp:
                tmp_path = tmp.name
            
            # 调用pdftotext
            cmd = ['pdftotext', '-layout', '-enc', 'UTF-8', str(file_path), tmp_path]
            result = subprocess.run(cmd, capture_output=True, text=True)
            
            if result.returncode != 0:
                logger.error(f"pdftotext执行失败: {result.stderr}")
                # 清理临时文件
                import os
                os.unlink(tmp_path)
                return None
            
            # 读取转换后的文本
            with open(tmp_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # 清理临时文件
            import os
            os.unlink(tmp_path)
            
            if not content or not content.strip():
                logger.warning(f"PDF文件无文本内容: {file_path.name}")
                return None
            
            logger.info(f"PDF解析完成 (pdftotext): {file_path.name}")
            return content.strip()
            
        except Exception as e:
            logger.error(f"pdftotext解析PDF文件失败 {file_path}: {e}")
            return None


class DocxParser(BaseParser):
    """Word文档解析器（.docx）"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.docx', '.doc']
        self._has_docx = self._check_docx_availability()
    
    def _check_docx_availability(self) -> bool:
        """检查python-docx是否可用"""
        try:
            import importlib.util
            spec = importlib.util.find_spec("docx")
            return spec is not None
        except Exception:
            logger.warning("python-docx库未安装，DOCX解析功能受限")
            return False
    
    def parse(self, file_path: Path) -> Optional[str]:
        """
        解析Word文档
        
        使用python-docx库提取文本内容（如果可用）
        否则尝试使用其他方法
        
        Args:
            file_path: Word文档路径
            
        Returns:
            文档文本内容，解析失败返回None
        """
        # 检查文件有效性
        if not self._is_valid_docx(file_path):
            logger.warning(f"无效的Word文档: {file_path.name}")
            return None
        
        if self._has_docx:
            return self._parse_with_docx(file_path)
        else:
            return self._parse_with_zip(file_path)
    
    def _parse_with_docx(self, file_path: Path) -> Optional[str]:
        """使用python-docx解析Word文档"""
        try:
            import docx
            
            doc = docx.Document(file_path)
            content_parts = []
            
            # 提取段落文本
            for para in doc.paragraphs:
                if para.text and para.text.strip():
                    content_parts.append(para.text.strip())
            
            # 提取表格文本
            for table in doc.tables:
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        if cell.text and cell.text.strip():
                            row_texts.append(cell.text.strip())
                    if row_texts:
                        content_parts.append(' | '.join(row_texts))
            
            if not content_parts:
                logger.warning(f"Word文档无文本内容: {file_path.name}")
                return None
            
            # 合并所有内容
            full_content = '\n\n'.join(content_parts)
            logger.info(f"Word文档解析完成 (python-docx): {file_path.name}")
            
            return full_content
            
        except Exception as e:
            logger.error(f"python-docx解析Word文档失败 {file_path}: {e}")
            # 尝试使用zip方法作为备选
            return self._parse_with_zip(file_path)
    
    def _parse_with_zip(self, file_path: Path) -> Optional[str]:
        """使用zipfile解析docx文件（docx本质是zip）"""
        try:
            import zipfile
            
            # docx文件实际上是zip压缩包，包含word/document.xml
            with zipfile.ZipFile(file_path, 'r') as docx_zip:
                # 查找文档内容
                if 'word/document.xml' not in docx_zip.namelist():
                    logger.error("无效的docx文件: 缺少word/document.xml")
                    return None
                
                # 读取document.xml
                with docx_zip.open('word/document.xml') as xml_file:
                    xml_content = xml_file.read().decode('utf-8')
                
                # 简单提取文本（实际应该解析XML命名空间）
                # 这里使用简单的方法：提取<w:t>标签内的文本
                import re
                text_matches = re.findall(r'<w:t[^>]*>([^<]+)</w:t>', xml_content)
                
                if not text_matches:
                    logger.warning(f"Word文档无文本内容: {file_path.name}")
                    return None
                
                # 合并文本，添加空格分隔
                full_content = ' '.join(text_matches)
                
                # 清理多余空格
                full_content = re.sub(r'\s+', ' ', full_content).strip()
                
                logger.info(f"Word文档解析完成 (zip方法): {file_path.name}")
                return full_content
                
        except zipfile.BadZipFile:
            logger.error(f"无效的zip文件（不是有效的docx）: {file_path.name}")
            return None
        except Exception as e:
            logger.error(f"zip方法解析Word文档失败 {file_path}: {e}")
            return None
    
    def _is_valid_docx(self, file_path: Path) -> bool:
        """检查是否为有效的docx文件"""
        try:
            # 检查文件扩展名
            if file_path.suffix.lower() not in ['.docx', '.doc']:
                return False
            
            # 检查文件大小（避免处理超大文件）
            file_size = file_path.stat().st_size
            if file_size > 100 * 1024 * 1024:  # 100MB限制
                logger.warning(f"Word文档过大 ({file_size/1024/1024:.1f}MB): {file_path.name}")
                return False
            
            # 检查文件是否可读
            if not os.access(file_path, os.R_OK):
                logger.warning(f"Word文档不可读: {file_path.name}")
                return False
            
            return True
        except Exception:
            return False


class CSVParser(BaseParser):
    """CSV文件解析器"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.csv']
        self._csv_available = True
    
    def parse(self, file_path: Path) -> Optional[str]:
        """
        解析CSV文件
        
        Args:
            file_path: CSV文件路径
            
        Returns:
            CSV文本内容，解析失败返回None
        """
        try:
            import csv
            
            content_parts = []
            encodings = ['utf-8', 'utf-8-sig', 'gbk', 'gb2312', 'latin-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        csv_reader = csv.reader(f)
                        content_parts = []
                        for row_num, row in enumerate(csv_reader, 1):
                            if row:
                                row_text = ' | '.join(str(cell).strip() for cell in row if str(cell).strip())
                                if row_text:
                                    content_parts.append(row_text)
                            if row_num % 100 == 0:
                                logger.debug(f"CSV解析进度: {row_num}行")
                    break  # 成功解析，退出编码循环
                except UnicodeDecodeError:
                    continue
            
            if not content_parts:
                logger.warning(f"CSV文件无文本内容: {file_path.name}")
                return None
            
            full_content = '\n'.join(content_parts)
            logger.info(f"CSV解析完成: {file_path.name}，字符数: {len(full_content)}")
            return full_content
        except ImportError:
            logger.error("csv模块不可用，无法解析CSV文件")
            return None
        except Exception as e:
            logger.error(f"解析CSV文件失败 {file_path.name}: {e}")
            return None


class HTMLParser(BaseParser):
    """HTML文件解析器"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.html', '.htm']
        try:
            from bs4 import BeautifulSoup
            self._has_bs4 = True
        except ImportError:
            self._has_bs4 = False
            logger.warning("bs4未安装，HTML解析将使用纯文本降级模式")
    
    def parse(self, file_path: Path) -> Optional[str]:
        """
        解析HTML文件
        
        Args:
            file_path: HTML文件路径
            
        Returns:
            纯文本内容，解析失败返回None
        """
        try:
            encodings = ['utf-8', 'utf-8-sig', 'gbk', 'gb2312', 'latin-1']
            raw_bytes = file_path.read_bytes()
            
            for encoding in encodings:
                try:
                    html_content = raw_bytes.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                logger.error(f"无法解码HTML文件: {file_path.name}")
                return None
            
            if self._has_bs4:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(html_content, 'html.parser')
                # 移除脚本和样式
                for tag in soup(['script', 'style', 'noscript', 'iframe']):
                    tag.decompose()
                text = soup.get_text(separator='\n', strip=True)
            else:
                # 降级：用正则去除HTML标签
                import re
                text = re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
                text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
                text = re.sub(r'<[^>]+>', ' ', text)
                text = re.sub(r'\s+', ' ', text).strip()
            
            if not text or len(text.strip()) < 10:
                logger.warning(f"HTML文件无有效文本: {file_path.name}")
                return None
            
            logger.info(f"HTML解析完成: {file_path.name}，字符数: {len(text)}")
            return text
        except Exception as e:
            logger.error(f"解析HTML文件失败 {file_path.name}: {e}")
            return None


class XlsxParser(BaseParser):
    """XLSX文件解析器（Excel）"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.xlsx', '.xls']
        try:
            import openpyxl
            self._has_openpyxl = True
        except ImportError:
            self._has_openpyxl = False
            logger.warning("openpyxl未安装，XLSX解析将使用ZIP降级模式")
    
    def parse(self, file_path: Path) -> Optional[str]:
        """
        解析XLSX文件
        
        Args:
            file_path: XLSX文件路径
            
        Returns:
            文本内容，解析失败返回None
        """
        try:
            content_parts = []
            
            if self._has_openpyxl:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    content_parts.append(f"工作表: {sheet_name}\n")
                    for row in sheet.iter_rows(values_only=True):
                        if any(str(cell).strip() for cell in row if cell is not None):
                            row_text = ' | '.join(str(cell).strip() for cell in row if cell is not None and str(cell).strip())
                            if row_text:
                                content_parts.append(row_text)
                wb.close()
            else:
                # 降级：ZIP方式读取sharedStrings.xml
                import zipfile
                import xml.etree.ElementTree as ET
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    if 'xl/sharedStrings.xml' in zip_ref.namelist():
                        with zip_ref.open('xl/sharedStrings.xml') as f:
                            tree = ET.parse(f)
                            root = tree.getroot()
                            ns = {'ns': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'}
                            for si in root.findall('.//ns:si', ns):
                                t = si.find('.//ns:t', ns)
                                if t is not None and t.text and t.text.strip():
                                    content_parts.append(t.text.strip())
            
            if not content_parts:
                logger.warning(f"XLSX文件无有效文本: {file_path.name}")
                return None
            
            full_content = '\n'.join(content_parts)
            logger.info(f"XLSX解析完成: {file_path.name}，字符数: {len(full_content)}")
            return full_content
        except Exception as e:
            logger.error(f"解析XLSX文件失败 {file_path.name}: {e}")
            return None


class PptxParser(BaseParser):
    """PPTX文件解析器（PowerPoint）"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.pptx', '.ppt']
        try:
            import pptx
            self._has_pptx = True
        except ImportError:
            self._has_pptx = False
            logger.warning("python-pptx未安装，PPTX解析将使用ZIP降级模式")
    
    def parse(self, file_path: Path) -> Optional[str]:
        """
        解析PPTX文件
        
        Args:
            file_path: PPTX文件路径
            
        Returns:
            文本内容，解析失败返回None
        """
        try:
            content_parts = []
            
            if self._has_pptx:
                import pptx
                prs = pptx.Presentation(file_path)
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    if slide_text:
                        content_parts.append(f"第{slide_num}页:")
                        content_parts.extend(slide_text)
            else:
                # 降级：ZIP方式读取slide*.xml
                import zipfile
                import xml.etree.ElementTree as ET
                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    slide_files = sorted(f for f in zip_ref.namelist() if 'ppt/slides/slide' in f and f.endswith('.xml'))
                    for slide_num, slide_file in enumerate(slide_files, 1):
                        try:
                            with zip_ref.open(slide_file) as f:
                                tree = ET.parse(f)
                                root = tree.getroot()
                                texts = []
                                for t in root.findall('.//a:t', ns):
                                    if t.text and t.text.strip():
                                        texts.append(t.text.strip())
                                if texts:
                                    content_parts.append(f"第{slide_num}页:")
                                    content_parts.extend(texts)
                        except Exception:
                            continue
            
            if not content_parts:
                logger.warning(f"PPTX文件无有效文本: {file_path.name}")
                return None
            
            full_content = '\n'.join(content_parts)
            logger.info(f"PPTX解析完成: {file_path.name}，字符数: {len(full_content)}")
            return full_content
        except Exception as e:
            logger.error(f"解析PPTX文件失败 {file_path.name}: {e}")
            return None


class ParserFactory:
    """解析器工厂类"""
    
    _parsers: List[BaseParser] = None
    
    @classmethod
    def get_parsers(cls) -> List[BaseParser]:
        """获取所有解析器实例"""
        if cls._parsers is None:
            cls._parsers = [
                TextParser(),
                MarkdownParser(),
                PDFParser(),
                DocxParser(),
                CSVParser(),
                HTMLParser(),
                XlsxParser(),
                PptxParser()
            ]
            # 所有解析器都保留（带降级逻辑，无依赖时也能工作）
            # 只过滤掉真正无法工作的解析器（通常不需要）
            pass
        return cls._parsers
    
    @classmethod
    def get_parser_for_file(cls, file_path: Path) -> Optional[BaseParser]:
        """
        根据文件扩展名获取合适的解析器
        
        Args:
            file_path: 文件路径对象
            
        Returns:
            匹配的解析器实例，无匹配返回None
        """
        for parser in cls.get_parsers():
            if parser.can_parse(file_path):
                return parser
        return None
    
    @classmethod
    def extract_title(cls, file_path: Path, content: Optional[str] = None) -> str:
        """
        提取文档标题
        
        Args:
            file_path: 文件路径对象
            content: 已解析的文本内容（可选）
            
        Returns:
            文档标题
        """
        parser = cls.get_parser_for_file(file_path)
        if parser and hasattr(parser, 'extract_title'):
            return parser.extract_title(file_path, content)
        return file_path.stem
    
    @classmethod
    def parse_file(cls, file_path: Path) -> Optional[str]:
        """
        解析单个文件
        
        Args:
            file_path: 文件路径对象
            
        Returns:
            文件文本内容，解析失败返回None
        """
        # 检查文件是否存在
        if not file_path.exists():
            logger.error(f"文件不存在: {file_path}")
            return None
        
        # 检查文件是否可读
        if not os.access(file_path, os.R_OK):
            logger.error(f"文件不可读: {file_path}")
            return None
        
        # 获取合适的解析器
        parser = cls.get_parser_for_file(file_path)
        if not parser:
            logger.warning(f"不支持的文件格式: {file_path.suffix}")
            return None
        
        # 解析文件
        logger.info(f"开始解析文件: {file_path.name}")
        content = parser.parse(file_path)
        
        if content:
            logger.info(f"文件解析成功: {file_path.name}，字符数: {len(content)}")
        else:
            logger.warning(f"文件解析失败: {file_path.name}")
        
        return content


# 全局函数接口
def parse_file(file_path: str) -> Optional[str]:
    """
    解析单个文件（全局函数接口）
    
    Args:
        file_path: 文件路径字符串
        
    Returns:
        文件文本内容，解析失败返回None
    """
    return ParserFactory.parse_file(Path(file_path))


def get_supported_extensions() -> List[str]:
    """获取所有支持的文件扩展名"""
    extensions = []
    for parser in ParserFactory.get_parsers():
        extensions.extend(parser.supported_extensions)
    return list(set(extensions))


if __name__ == "__main__":
    # 测试代码
    import sys
    
    if len(sys.argv) > 1:
        file_path = Path(sys.argv[1])
        content = parse_file(file_path)
        if content:
            print(f"解析成功，字符数: {len(content)}")
            print("前500字符:")
            print(content[:500])
        else:
            print("解析失败")
    else:
        print("支持的扩展名:", get_supported_extensions())